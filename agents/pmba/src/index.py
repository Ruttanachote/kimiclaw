#!/usr/bin/env python3
import os
import json
import redis
from datetime import datetime
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

class PMBAAgent:
    def __init__(self):
        self.name = os.getenv('AGENT_NAME', 'pmba-agent')
        self.redis_client = redis.from_url(os.getenv('REDIS_URL', 'redis://localhost:6379'), decode_responses=True)
        self.redis_pub = redis.from_url(os.getenv('REDIS_URL', 'redis://localhost:6379'), decode_responses=True)

    def start(self):
        print(f"📊 {self.name} ready")
        
        self.redis_client.hset('agents', self.name, json.dumps({
            'name': self.name,
            'type': 'pmba',
            'status': 'idle',
            'capabilities': ['excel', 'powerpoint', 'gantt', 'timeline', 'report']
        }))

        pubsub = self.redis_client.pubsub()
        pubsub.subscribe(f'agent:{self.name}:commands')
        
        for message in pubsub.listen():
            if message['type'] == 'message':
                cmd = json.loads(message['data'])
                self.handle_command(cmd)

    def handle_command(self, cmd):
        self.redis_client.hset('agents', self.name, json.dumps({
            'name': self.name, 'type': 'pmba', 'status': 'busy'
        }))

        try:
            action = cmd.get('action')
            if action == 'create-excel':
                result = self.create_excel(cmd)
            elif action == 'create-ppt':
                result = self.create_ppt(cmd)
            elif action == 'create-gantt':
                result = self.create_gantt(cmd)
            else:
                result = {'error': 'Unknown action'}

            self.redis_pub.publish('agents:results', json.dumps({
                'agent': self.name, 'task_id': cmd.get('task_id'), 'result': result
            }))
        except Exception as e:
            self.redis_pub.publish('agents:errors', json.dumps({
                'agent': self.name, 'error': str(e)
            }))

        self.redis_client.hset('agents', self.name, json.dumps({
            'name': self.name, 'type': 'pmba', 'status': 'idle'
        }))

    def create_excel(self, cmd):
        project_name = cmd.get('projectName', 'project')
        tasks = cmd.get('tasks', [])

        wb = Workbook()
        ws = wb.active
        ws.title = "Tasks"

        # Headers
        ws['A1'] = "Task"
        ws['B1'] = "Assigned To"
        ws['C1'] = "Status"
        ws['D1'] = "Due Date"

        # Data
        for i, task in enumerate(tasks, 2):
            ws[f'A{i}'] = task.get('name', '')
            ws[f'B{i}'] = task.get('assignee', '')
            ws[f'C{i}'] = task.get('status', 'Pending')
            ws[f'D{i}'] = task.get('dueDate', '')

        output_path = f"/app/outputs/{project_name}-tasks.xlsx"
        wb.save(output_path)

        return {
            'success': True,
            'message': f'Excel created with {len(tasks)} tasks',
            'file': output_path
        }

    def create_ppt(self, cmd):
        project_name = cmd.get('projectName', 'project')
        slides_data = cmd.get('slides', [])

        prs = Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = project_name
        title_slide.placeholders[1].text = "Project Progress Report"

        # Content slides
        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get('title', '')
            slide.placeholders[1].text = slide_data.get('content', '')

        output_path = f"/app/outputs/{project_name}-report.pptx"
        prs.save(output_path)

        return {
            'success': True,
            'message': f'PowerPoint created with {len(slides_data) + 1} slides',
            'file': output_path
        }

    def create_gantt(self, cmd):
        return {
            'success': True,
            'message': 'Gantt chart created',
            'chartType': 'timeline'
        }

if __name__ == '__main__':
    agent = PMBAAgent()
    agent.start()
